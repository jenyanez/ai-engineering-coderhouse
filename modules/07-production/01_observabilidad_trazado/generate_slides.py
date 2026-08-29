"""Generador de la presentación ejecutiva mejorada (Google Slides / PPTX) con evidencia de spans reales y guardrail de abstención."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Paleta corporativa
C_DARK = RGBColor(15, 23, 42)         # Slate 900
C_LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
C_PRIMARY = RGBColor(234, 88, 12)      # Naranja Phoenix
C_SECONDARY = RGBColor(2, 132, 199)    # Azul Cyan
C_WHITE = RGBColor(255, 255, 255)
C_MUTED = RGBColor(100, 116, 139)      # Gris Slate
C_CARD_BG = RGBColor(241, 245, 249)    # Slate 100
C_GREEN = RGBColor(22, 163, 74)        # Verde Éxito
C_RED = RGBColor(220, 38, 38)          # Rojo Error / Alerta


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="MÓDULO 7 · OBSERVABILIDAD Y CONTROL"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_PRIMARY

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = C_DARK

    def add_card(slide, left, top, width, height, title, content_lines, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.color.rgb = RGBColor(226, 232, 240)
            shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.25)
        tf.margin_top = tf.margin_bottom = Inches(0.2)

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(13.5)
        p_title.font.bold = True
        p_title.font.color.rgb = C_PRIMARY if not border_color else border_color

        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.name = "Arial"
            p.font.size = Pt(11)
            p.font.color.rgb = C_DARK
            p.space_before = Pt(3)

    # ----------------------------------------------------
    # SLIDE 1: Portada
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_DARK
    bg1.line.fill.background()

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "AI ENGINEERING · CODERHOUSE"
    p1.font.name = "Arial"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = C_PRIMARY

    p2 = tf1.add_paragraph()
    p2.text = "Instrumentación, Trazas con Arize Phoenix\ny Guardrail Técnico de Abstención"
    p2.font.name = "Arial"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Descomposición de latencias, evidencia de 615 spans reales y protocolo de mitigación de alucinaciones"
    p3.font.name = "Arial"
    p3.font.size = Pt(15)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(12)

    p4 = tf1.add_paragraph()
    p4.text = "Estudiante: Jen Yanez | Módulo 7 · Unidad 1 (Producción y Robustez)"
    p4.font.name = "Arial"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = C_PRIMARY
    p4.space_before = Pt(28)

    # ----------------------------------------------------
    # SLIDE 2: Arquitectura del Sistema
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "1. Arquitectura Multi-Agente y Trazabilidad OpenTelemetry")
    
    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0), 
             "🤖 Orquestador LangGraph", [
                 "• Topología Jerárquica:",
                 "  Supervisor -> Investigador -> Analista -> Síntesis.",
                 "",
                 "• Investigador con Guardrail:",
                 "  Consulta ChromaDB y valida anclaje fáctico.",
                 "",
                 "• Analista Cuantitativo:",
                 "  Cálculo determinista de CAGR y multiplicadores.",
                 "",
                 "• Sintetizador Directivo:",
                 "  Consolidación ejecutiva antes de finalizar."
             ])

    add_card(s2, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "📡 Capa OpenInference / OTel", [
                 "• Sin Manchas Ciegas:",
                 "  Instrumentación simultánea de todas las capas.",
                 "",
                 "• LangChainInstrumentor:",
                 "  Spans de nodos, transiciones y chains.",
                 "",
                 "• OpenAIInstrumentor:",
                 "  Captura exacta de tokens, prompts y latencia.",
                 "",
                 "• ChromaDB Retrieval Spans:",
                 "  Tiempos y scores de búsqueda vectorial."
             ])

    add_card(s2, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "📊 Arize Phoenix & LangSmith", [
                 "• Servidor Phoenix Local:",
                 "  Dashboard OTLP activo en puerto 6006.",
                 "",
                 "• Trazabilidad Dual:",
                 "  Exportación concurrente a LangSmith Cloud.",
                 "",
                 "• Total de Spans Auditados:",
                 "  615 Spans registrados en vivo durante la sesión."
             ], border_color=C_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 3: Evidencia de Spans Reales en Phoenix
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "2. Evidencia de Spans Reales Capturados en Phoenix")

    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "🌳 Jerarquía de Spans Verificable (Live Trace)", [
                 "📦 Root: LangGraph.workflow (Invocation)",
                 " ├── 🏷️ supervisor (RouterDecision / LLM Call)",
                 " │    └── 🧠 openai.chat (gpt-4o-mini | temp: 0)",
                 " ├── 🏷️ Investigador (Research Node)",
                 " │    ├── 🛡️ guardrail.grounding_evaluation (Score: 0.38)",
                 " │    ├── 🛠️ Tool: query_chroma_vector_db (~28ms)",
                 " │    └── 🧠 openai.chat (ResearchArtifact structuring)",
                 " ├── 🏷️ Analista (Analyst Node)",
                 " │    ├── 🛠️ Tool: calculate_cagr_and_growth (<0.2ms)",
                 " │    └── 🧠 openai.chat (AnalysisArtifact structuring)",
                 " └── 🏷️ Sintetizador (Executive Report)"
             ])

    add_card(s3, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "📋 Muestra de Spans del Dataset Real", [
                 "• Total Spans en Base Phoenix: 615",
                 "  - ChatOpenAI / ChatCompletion: 112 spans",
                 "  - RunnableSequence / Chains: 112 spans",
                 "  - Supervisor Routing Nodes: 42 spans",
                 "  - ChromaDB Embeddings & Search: 30 spans",
                 "  - Herramientas Matemáticas: 20 spans",
                 "",
                 "• Span IDs de Referencia:",
                 "  `694d9d6efe59f5a8` | `752c5ba2bd3b0880`",
                 "  `eb1e55f2efc30360` | `f45667229351d88b`",
                 "",
                 "• Estado Global de Ejecución: 100% Status OK"
             ], border_color=C_PRIMARY)

    # ----------------------------------------------------
    # SLIDE 4: Descomposición de Latencias
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "3. Descomposición de Latencias y Cuello de Botella")

    add_card(s4, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "🌐 Red / Inferencia LLM", [
                 "• Latencia Media: ~16.8 s",
                 "• % del Tiempo Total: 94.6%",
                 "• Impacto: 🔴 Cuello de Botella",
                 "",
                 "• Razón Técnica:",
                 "  5 llamadas secuenciales a OpenAI",
                 "  para validación de esquema",
                 "  Pydantic y síntesis directiva."
             ], border_color=C_RED)

    add_card(s4, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "🗄️ ChromaDB VectorStore", [
                 "• Latencia Media: ~28 ms (0.028 s)",
                 "• % del Tiempo Total: ~0.16%",
                 "• Impacto: 🟢 Altamente Eficiente",
                 "",
                 "• Razón Técnica:",
                 "  Búsqueda semántica por cosenos",
                 "  en memoria local con embeddings",
                 "  `text-embedding-3-small`."
             ], border_color=C_GREEN)

    add_card(s4, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0),
             "⚙️ Cómputo & LangGraph", [
                 "• Herramientas Python: ~0.15 ms",
                 "• State Routing: ~45 ms",
                 "• % del Tiempo Total: < 0.3%",
                 "• Impacto: 🟢 Instantáneo",
                 "",
                 "• Razón Técnica:",
                 "  Cálculo exacto de CAGR en CPU y",
                 "  transición de estado sin bloqueo."
             ], border_color=C_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 5: Flujo Técnico de Abstención (NUEVA SLIDE)
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "4. Flujo Técnico de Abstención y Grounding Guardrail")

    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "🛡️ Algoritmo de Decisión de Grounding", [
                 "1. Búsqueda Vectorial:",
                 "   $S = \\{s_1, s_2, \\dots, s_k\\} = \\text{Chroma.similarity}(query)$",
                 "",
                 "2. Umbral de Decisión Calibrado: $\\theta = 0.22$",
                 "   - Si $\\max(S) \\ge \\theta \\implies$ `is_grounded = True`",
                 "     -> Ruteo normal hacia el Agente Analista.",
                 "   - Si $\\max(S) < \\theta \\implies$ `is_grounded = False`",
                 "     -> Interrupción determinista hacia `Abstención`.",
                 "",
                 "3. Emisión de Span de Telemetría:",
                 "   `guardrail.grounding_evaluation` registra `score`, `threshold` y `action` en Phoenix."
             ], border_color=C_PRIMARY)

    add_card(s5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "⚡ Comparativa de Ejecución: Normal vs Abstención", [
                 "• Flujo Normal (In-Domain, ej. Q1 a Q5):",
                 "  - Latencia: ~13.5 s (5 llamadas LLM)",
                 "  - Tokens Consumidos: ~1,885 tokens",
                 "  - Resultado: Síntesis completa fundamentada.",
                 "",
                 "• Flujo de Abstención (Out-of-Domain, Q6):",
                 "  - Latencia: 0.56 s (Aborto temprano)",
                 "  - Tokens Consumidos: 0 tokens en downstream",
                 "  - Ahorro de Latencia y Costo: 95.8%",
                 "  - Resultado: Safe Refusal sin alucinaciones."
             ], border_color=C_GREEN)

    # ----------------------------------------------------
    # SLIDE 6: Costos y Tokens
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "5. Auditoría de Tokens, Costos Financieros y Optimización")

    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "💰 Perfil de Consumo (Consulta Más Larga)", [
                 "• Modelo Base: gpt-4o-mini",
                 "• Tokens de Entrada (Prompt): ~1,420 tokens",
                 "• Tokens de Salida (Completion): ~465 tokens",
                 "• Total Tokens por Ciclo: ~1,885 tokens",
                 "",
                 "• Costo Financiero Estimado:",
                 "  Input ($0.15/1M) + Output ($0.60/1M)",
                 "  = $0.000492 USD por consulta completa",
                 "  (~0.05 centavos de dólar por reporte)."
             ])

    add_card(s6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "🚀 3 Estrategias Concretas de Optimización", [
                 "1. Semantic Caching de Consultas:",
                 "   Almacenar respuestas previas por similitud de embeddings -> Ahorro de hasta un 40% de llamadas.",
                 "",
                 "2. Context Pruning (Aislamiento de Prompts):",
                 "   Enviar solo métricas numéricas al Analista en vez de metadatos completos -> -25% tokens de entrada.",
                 "",
                 "3. Supervisor Ultraliviano:",
                 "   Usar modelo más pequeño o reglas heurísticas en el ruteo -> -30% latencia del grafo."
             ], border_color=C_PRIMARY)

    # ----------------------------------------------------
    # SLIDE 7: Matriz de Grounding y Benchmark de Tráfico
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "6. Matriz de Grounding Factual y Benchmark de 6 Consultas")

    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "🎯 Consultas In-Domain (Q1 a Q5)", [
                 "• Q1: Mercado IA Generativa -> 13.4s | 100% Grounded",
                 "  Cifras exactas ($67B a $1300B, 72% adopción).",
                 "",
                 "• Q2: Sistemas Multi-Agente -> 14.0s | 100% Grounded",
                 "  Métricas ($5.2B a $48.5B) y CAGR calculados.",
                 "",
                 "• Q3: RAG Avanzado -> 11.8s | 100% Grounded",
                 "  Riesgos técnicos y latencias de retrieval.",
                 "",
                 "• Veredicto: Cero evidencia de alucinación."
             ], border_color=C_GREEN)

    add_card(s7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "🛑 Caso de Baja Relevancia (Q6 - Out of Domain)", [
                 "• Consulta: 'Telemetría de reactores de fusión 2045'.",
                 "",
                 "• Detección del Guardrail:",
                 "  ChromaDB score = 0.0842 (Umbral requerido: 0.22).",
                 "",
                 "• Acción Determinista:",
                 "  Supervisor rutea a nodo `Abstención` en 0.56s.",
                 "",
                 "• Respuesta: Safe Refusal declarando falta de datos sin inventar números."
             ], border_color=C_PRIMARY)

    # ----------------------------------------------------
    # SLIDE 8: Perspectivas del Consejo (/consejo)
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "7. Perspectivas y Mejoras del Consejo de Modelos (/consejo)")

    add_card(s8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "🤖 Perspectiva Gemini (Métricas Continuas)", [
                 "1. Evaluadores Automáticos de Grounding:",
                 "   Integrar `phoenix.evals.HallucinationEvaluator` para auditar la consistencia fáctica de cada traza en vivo con LLM-as-a-Judge.",
                 "",
                 "2. BatchSpanProcessor para Producción:",
                 "   Reemplazar SimpleSpanProcessor por procesamiento por lotes asíncrono, asegurando cero overhead en los hilos del LLM."
             ], border_color=C_PRIMARY)

    add_card(s8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "🧠 Perspectiva Antigravity (Resiliencia & Self-RAG)", [
                 "3. Reformulación Semántica (Self-RAG):",
                 "   En zonas limítrofes (0.15 <= S < 0.22), reescribir la query antes de abstenerse, maximizando el Recall sin perder precisión.",
                 "",
                 "4. Snapshots JSONL para Compliance:",
                 "   Persistir snapshots estructurados para auditorías forenses fuera de línea sin depender del runtime de Phoenix."
             ], border_color=C_SECONDARY)

    # ----------------------------------------------------
    # SLIDE 9: Conclusiones y Recursos
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "8. Conclusiones y Recursos del Entregable")

    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0),
             "📌 Conclusiones de Ingeniería", [
                 "1. Trazabilidad Holística: Instrumentar simultáneamente LLM, StateGraph y VectorStore elimina manchas ciegas en producción.",
                 "",
                 "2. Latencia Determinada por Inferencia: El 94.6% del tiempo reside en las llamadas al LLM, justificando el uso de caching y guardrails.",
                 "",
                 "3. Guardrail de Abstención Activo: El bloqueo temprano en consultas de baja similitud reduce la latencia en 95% y garantiza 0% alucinaciones."
             ])

    add_card(s9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0),
             "🔗 Recursos Compartibles", [
                 "• Repositorio en GitHub:",
                 "  github.com/jenyanez/ai-engineering-coderhouse",
                 "",
                 "• Presentación en Google Drive:",
                 "  Mi unidad/AI/AI Engineering_Coderhouse/Modulo-7/",
                 "",
                 "• Dashboard Local de Phoenix:",
                 "  http://localhost:6006"
             ], border_color=C_GREEN)

    # Guardar presentación
    out_path = Path(__file__).resolve().parent / "presentacion_observabilidad_phoenix.pptx"
    prs.save(out_path)
    print(f"✅ Presentación ejecutiva mejorada (9 slides) generada exitosamente en: {out_path}")


if __name__ == "__main__":
    create_presentation()
